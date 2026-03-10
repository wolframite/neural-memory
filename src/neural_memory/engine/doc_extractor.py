"""Multi-format document extractor for doc-to-brain training pipeline.

Converts supported file formats to markdown text, which is then fed into
the existing chunk_markdown() pipeline. Each extractor handles one format
and outputs clean markdown suitable for heading-based chunking.

Supported formats:
- Text passthrough: .md, .mdx, .txt, .rst
- Rich documents: .pdf, .docx, .pptx, .html, .htm
- Structured data: .json, .xlsx, .csv
"""

from __future__ import annotations

import csv
import io
import json
import logging
from collections.abc import Callable
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Maximum file size to process (2GB)
_MAX_FILE_SIZE = 2 * 1024 * 1024 * 1024

# Threshold for page-by-page PDF extraction (50MB)
_PDF_STREAMING_THRESHOLD = 50 * 1024 * 1024

# Maximum rows for tabular data (CSV/XLSX)
_MAX_TABLE_ROWS = 10_000

# All extensions supported by the extractor
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".md",
        ".mdx",
        ".txt",
        ".rst",
        ".pdf",
        ".docx",
        ".pptx",
        ".html",
        ".htm",
        ".json",
        ".xlsx",
        ".csv",
    }
)

# Extensions that need optional dependencies
_RICH_DOC_EXTENSIONS: frozenset[str] = frozenset({".pdf", ".docx", ".pptx", ".html", ".htm"})
_SPREADSHEET_EXTENSIONS: frozenset[str] = frozenset({".xlsx"})


class ExtractionError(Exception):
    """Raised when document extraction fails."""


def extract_to_markdown(file_path: Path) -> str:
    """Extract any supported file to markdown string.

    Args:
        file_path: Path to the file to extract.

    Returns:
        Markdown text suitable for chunk_markdown().

    Raises:
        ExtractionError: If extraction fails or format is unsupported.
        FileNotFoundError: If file does not exist.
    """
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    if not file_path.is_file():
        raise ExtractionError(f"Not a file: {file_path}")

    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ExtractionError(
            f"Unsupported format: {suffix}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    # Check file size
    file_size = file_path.stat().st_size
    if file_size > _MAX_FILE_SIZE:
        raise ExtractionError(
            f"File too large: {file_size / 1024 / 1024:.1f}MB (max {_MAX_FILE_SIZE / 1024 / 1024:.0f}MB)"
        )

    extractor = _EXTRACTOR_MAP.get(suffix)
    if extractor is None:
        raise ExtractionError(f"No extractor for: {suffix}")

    return extractor(file_path)


def get_missing_dependencies(extensions: frozenset[str] | set[str]) -> list[str]:
    """Check which optional dependencies are missing for the given extensions.

    Returns:
        List of missing package names (empty if all available).
    """
    missing: list[str] = []

    if extensions & {".pdf"}:
        try:
            import pymupdf4llm  # noqa: F401
        except ImportError:
            missing.append("pymupdf4llm")

    if extensions & {".docx"}:
        try:
            import docx  # noqa: F401
        except ImportError:
            missing.append("python-docx")

    if extensions & {".pptx"}:
        try:
            import pptx  # noqa: F401
        except ImportError:
            missing.append("python-pptx")

    if extensions & {".html", ".htm"}:
        try:
            import bs4  # noqa: F401
            import markdownify  # noqa: F401
        except ImportError:
            missing.append("beautifulsoup4 markdownify")

    if extensions & {".xlsx"}:
        try:
            import openpyxl  # noqa: F401
        except ImportError:
            missing.append("openpyxl")

    return missing


# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------


def _extract_text(file_path: Path) -> str:
    """Passthrough for markdown/text files."""
    return file_path.read_text(encoding="utf-8")


def _extract_pdf(file_path: Path) -> str:
    """Extract PDF to markdown using pymupdf4llm.

    For files larger than _PDF_STREAMING_THRESHOLD (50MB), uses page-by-page
    extraction to avoid loading the entire document into memory at once.
    """
    try:
        import pymupdf4llm
    except ImportError:
        raise ExtractionError(
            "PDF extraction requires pymupdf4llm. Install with: pip install neural-memory[extract]"
        )

    file_size = file_path.stat().st_size
    if file_size <= _PDF_STREAMING_THRESHOLD:
        # Small PDF — extract all at once (fast)
        try:
            md_text: str = pymupdf4llm.to_markdown(str(file_path))
            return md_text
        except Exception as exc:
            raise ExtractionError(f"PDF extraction failed: {exc}")

    # Large PDF — page-by-page extraction to limit RAM usage
    return _extract_pdf_paged(file_path)


def _extract_pdf_paged(file_path: Path) -> str:
    """Page-by-page PDF extraction for large files (>50MB).

    Opens the PDF once with pymupdf, then extracts pages in batches
    to keep memory bounded. Each batch is converted to markdown and
    appended to the result.
    """
    try:
        import pymupdf  # pymupdf4llm depends on pymupdf
        import pymupdf4llm
    except ImportError:
        raise ExtractionError(
            "PDF extraction requires pymupdf4llm. Install with: pip install neural-memory[extract]"
        )

    page_batch = 20  # pages per batch to balance speed vs memory

    try:
        doc = pymupdf.open(str(file_path))
    except Exception as exc:
        raise ExtractionError(f"Failed to open PDF: {exc}")

    total_pages = len(doc)
    doc.close()

    logger.info(
        "Large PDF detected (%.1fMB, %d pages) — using paged extraction",
        file_path.stat().st_size / (1024 * 1024),
        total_pages,
    )

    parts: list[str] = []
    for start in range(0, total_pages, page_batch):
        end = min(start + page_batch, total_pages) - 1  # pymupdf4llm uses inclusive end
        try:
            batch_md: str = pymupdf4llm.to_markdown(
                str(file_path),
                pages=list(range(start, end + 1)),
            )
            parts.append(batch_md)
        except Exception as exc:
            logger.warning(
                "Failed to extract pages %d-%d from %s: %s",
                start,
                end,
                file_path.name,
                exc,
            )
            # Continue with remaining pages instead of aborting
            continue

    if not parts:
        raise ExtractionError("All page batches failed during paged extraction")

    return "\n\n".join(parts)


def _extract_docx(file_path: Path) -> str:
    """Extract DOCX to markdown using python-docx."""
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # noqa: F401
    except ImportError:
        raise ExtractionError(
            "DOCX extraction requires python-docx. Install with: pip install neural-memory[extract]"
        )

    try:
        doc = Document(str(file_path))
    except Exception as exc:
        raise ExtractionError(f"DOCX extraction failed: {exc}")

    lines: list[str] = []
    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        # Map Word heading styles to markdown headings
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.split()[-1])
                level = min(level, 6)
            except (ValueError, IndexError):
                level = 1
            lines.append(f"{'#' * level} {text}")
        elif style_name == "Title":
            lines.append(f"# {text}")
        elif style_name.startswith("List"):
            lines.append(f"- {text}")
        else:
            lines.append(text)

    # Extract tables
    for table in doc.tables:
        lines.append("")
        _docx_table_to_md(table, lines)

    return "\n".join(lines)


def _docx_table_to_md(table: Any, lines: list[str]) -> None:
    """Convert a python-docx table to markdown table."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return

    # Header row
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")

    # Data rows
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    lines.append("")


def _extract_pptx(file_path: Path) -> str:
    """Extract PPTX to markdown using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ExtractionError(
            "PPTX extraction requires python-pptx. Install with: pip install neural-memory[extract]"
        )

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        raise ExtractionError(f"PPTX extraction failed: {exc}")

    lines: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_title = ""

        # Find title shape
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_title = slide.shapes.title.text.strip()

        if slide_title:
            lines.append(f"## Slide {slide_idx}: {slide_title}")
        else:
            lines.append(f"## Slide {slide_idx}")
        lines.append("")

        # Extract text from all shapes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text and text != slide_title:
                    # Check if it's a bullet point (has indentation level)
                    indent = paragraph.level if paragraph.level else 0
                    if indent > 0:
                        lines.append(f"{'  ' * indent}- {text}")
                    else:
                        lines.append(text)

        # Extract tables
        for shape in slide.shapes:
            if shape.has_table:
                _pptx_table_to_md(shape.table, lines)

        lines.append("")

    return "\n".join(lines)


def _pptx_table_to_md(table: Any, lines: list[str]) -> None:
    """Convert a python-pptx table to markdown table."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return

    lines.append("")
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    lines.append("")


def _extract_html(file_path: Path) -> str:
    """Extract HTML to markdown using beautifulsoup4 + markdownify."""
    try:
        from bs4 import BeautifulSoup
        from markdownify import markdownify
    except ImportError:
        raise ExtractionError(
            "HTML extraction requires beautifulsoup4 and markdownify. "
            "Install with: pip install neural-memory[extract]"
        )

    try:
        raw_html = file_path.read_text(encoding="utf-8")
        soup = BeautifulSoup(raw_html, "html.parser")

        # Remove script, style, nav, footer elements
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        # Find main content area or fall back to body
        main = soup.find("main") or soup.find("article") or soup.find("body") or soup
        md_text: str = markdownify(str(main), heading_style="ATX", strip=["img"])
        return md_text
    except Exception as exc:
        raise ExtractionError(f"HTML extraction failed: {exc}")


def _extract_json(file_path: Path) -> str:
    """Extract JSON to structured markdown."""
    try:
        raw = file_path.read_text(encoding="utf-8")
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ExtractionError(f"JSON parse error: {exc}")
    except (OSError, UnicodeDecodeError) as exc:
        raise ExtractionError(f"Failed to read JSON file: {exc}")

    lines: list[str] = [f"# {file_path.stem}", ""]

    if isinstance(data, list):
        _json_array_to_md(data, lines, file_path.stem)
    elif isinstance(data, dict):
        _json_object_to_md(data, lines, depth=2)
    else:
        lines.append(f"Value: `{data}`")

    return "\n".join(lines)


def _json_object_to_md(obj: dict[str, Any], lines: list[str], depth: int = 2) -> None:
    """Convert a JSON object to markdown sections."""
    capped_depth = min(depth, 6)
    for key, value in obj.items():
        if isinstance(value, dict):
            lines.append(f"{'#' * capped_depth} {key}")
            lines.append("")
            _json_object_to_md(value, lines, depth=depth + 1)
        elif isinstance(value, list):
            lines.append(f"{'#' * capped_depth} {key}")
            lines.append("")
            _json_array_to_md(value, lines, key)
        else:
            lines.append(f"- **{key}**: {value}")

    lines.append("")


def _json_array_to_md(arr: list[Any], lines: list[str], context: str = "") -> None:
    """Convert a JSON array to markdown."""
    if not arr:
        lines.append("_(empty array)_")
        return

    # Check if array of objects → table
    if all(isinstance(item, dict) for item in arr):
        _dicts_to_md_table(arr, lines)
    else:
        for item in arr[:_MAX_TABLE_ROWS]:
            if isinstance(item, dict):
                for k, v in item.items():
                    lines.append(f"- **{k}**: {v}")
                lines.append("")
            else:
                lines.append(f"- {item}")

    if len(arr) > _MAX_TABLE_ROWS:
        lines.append(f"\n_(truncated: showing {_MAX_TABLE_ROWS} of {len(arr)} items)_")

    lines.append("")


def _dicts_to_md_table(items: list[dict[str, Any]], lines: list[str]) -> None:
    """Convert a list of flat dicts to a markdown table."""
    # Collect all keys across items
    all_keys: list[str] = []
    seen: set[str] = set()
    for item in items[:_MAX_TABLE_ROWS]:
        for key in item:
            if key not in seen:
                all_keys.append(key)
                seen.add(key)

    if not all_keys:
        return

    # Header
    lines.append("| " + " | ".join(all_keys) + " |")
    lines.append("| " + " | ".join("---" for _ in all_keys) + " |")

    # Rows
    for item in items[:_MAX_TABLE_ROWS]:
        cells = [str(item.get(k, "")).replace("|", "\\|").replace("\n", " ") for k in all_keys]
        lines.append("| " + " | ".join(cells) + " |")

    if len(items) > _MAX_TABLE_ROWS:
        lines.append(f"\n_(truncated: showing {_MAX_TABLE_ROWS} of {len(items)} rows)_")

    lines.append("")


def _extract_xlsx(file_path: Path) -> str:
    """Extract XLSX to markdown tables using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ExtractionError(
            "XLSX extraction requires openpyxl. Install with: pip install neural-memory[extract]"
        )

    try:
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
    except Exception as exc:
        raise ExtractionError(f"XLSX extraction failed: {exc}")

    lines: list[str] = [f"# {file_path.stem}", ""]

    try:
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            lines.append(f"## {sheet}")
            lines.append("")

            rows: list[list[str]] = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count >= _MAX_TABLE_ROWS:
                    lines.append(f"\n_(truncated at {_MAX_TABLE_ROWS} rows)_")
                    break
                cells = [str(cell) if cell is not None else "" for cell in row]
                # Skip entirely empty rows
                if any(c.strip() for c in cells):
                    rows.append(cells)
                    row_count += 1

            if rows:
                # First row as header
                lines.append("| " + " | ".join(rows[0]) + " |")
                lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                for row in rows[1:]:
                    # Pad row to match header length
                    padded = row + [""] * max(0, len(rows[0]) - len(row))
                    cells_escaped = [c.replace("|", "\\|").replace("\n", " ") for c in padded]
                    lines.append("| " + " | ".join(cells_escaped) + " |")

            lines.append("")
    finally:
        wb.close()

    return "\n".join(lines)


def _extract_csv(file_path: Path) -> str:
    """Extract CSV to markdown table."""
    try:
        raw = file_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError) as exc:
        raise ExtractionError(f"Failed to read CSV file: {exc}")

    lines: list[str] = [f"# {file_path.stem}", ""]

    reader = csv.reader(io.StringIO(raw))
    rows: list[list[str]] = []

    for row_count, row in enumerate(reader):
        if row_count >= _MAX_TABLE_ROWS:
            lines.append(f"\n_(truncated at {_MAX_TABLE_ROWS} rows)_")
            break
        if any(cell.strip() for cell in row):
            rows.append(row)

    if rows:
        # Normalize column count
        max_cols = max(len(r) for r in rows)
        lines.append("| " + " | ".join(rows[0] + [""] * (max_cols - len(rows[0]))) + " |")
        lines.append("| " + " | ".join("---" for _ in range(max_cols)) + " |")
        for row in rows[1:]:
            padded = row + [""] * (max_cols - len(row))
            cells_escaped = [c.replace("|", "\\|").replace("\n", " ") for c in padded]
            lines.append("| " + " | ".join(cells_escaped) + " |")

    lines.append("")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Extractor dispatch map
# ---------------------------------------------------------------------------

_EXTRACTOR_MAP: dict[str, Callable[[Path], str]] = {
    ".md": _extract_text,
    ".mdx": _extract_text,
    ".txt": _extract_text,
    ".rst": _extract_text,
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".json": _extract_json,
    ".xlsx": _extract_xlsx,
    ".csv": _extract_csv,
}
